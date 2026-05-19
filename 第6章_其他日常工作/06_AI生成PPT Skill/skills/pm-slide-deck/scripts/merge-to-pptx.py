#!/usr/bin/env python3
"""pm-slide-deck: 合并幻灯片图片为 PPTX

用法：
    python merge-to-pptx.py <幻灯片目录> [--output 文件名.pptx]

依赖：
    pip install python-pptx
"""

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def find_slide_images(dir_path):
    d = Path(dir_path)
    if not d.exists():
        print(f"目录不存在：{dir_path}", file=sys.stderr)
        sys.exit(1)

    slides = []
    for f in sorted(d.iterdir()):
        m = re.match(r"^(\d+)-slide-.*\.(png|jpg|jpeg)$", f.name, re.IGNORECASE)
        if m:
            slides.append({"filename": f.name, "path": str(f), "index": int(m.group(1))})

    slides.sort(key=lambda s: s["index"])

    if not slides:
        print(f"未找到幻灯片图片：{dir_path}", file=sys.stderr)
        print("期望格式：01-slide-*.png, 02-slide-*.png 等", file=sys.stderr)
        sys.exit(1)

    return slides


def find_prompts_for_slides(slides_dir, slides):
    prompts_dir = Path(slides_dir) / "prompts"
    if not prompts_dir.exists():
        return {}

    prompt_map = {}
    for slide in slides:
        base_name = Path(slide["filename"]).stem
        prompt_file = prompts_dir / f"{base_name}.md"
        if prompt_file.exists():
            prompt_map[slide["index"]] = prompt_file.read_text(encoding="utf-8")
    return prompt_map


def find_base_prompt():
    base_prompt_path = Path(__file__).resolve().parent.parent / "references" / "base-prompt.md"
    if base_prompt_path.exists():
        return base_prompt_path.read_text(encoding="utf-8")
    return None


def create_pptx(slides, output_path, slides_dir):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_prompt = find_base_prompt()
    prompt_map = find_prompts_for_slides(slides_dir, slides)
    notes_count = 0

    blank_layout = prs.slide_layouts[6]

    for slide in slides:
        s = prs.slides.add_slide(blank_layout)
        s.shapes.add_picture(
            slide["path"],
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height,
        )

        if slide["index"] in prompt_map:
            slide_prompt = prompt_map[slide["index"]]
            full_notes = f"{base_prompt}\n\n---\n\n{slide_prompt}" if base_prompt else slide_prompt
            s.notes_slide.notes_text_frame.text = full_notes
            notes_count += 1

        prompt_tag = "（含备注）" if slide["index"] in prompt_map else ""
        print(f"已添加：{slide['filename']}{prompt_tag}")

    prs.save(output_path)
    print(f"\n已创建：{output_path}")
    print(f"总页数：{len(slides)}")
    if notes_count > 0:
        print(f"含备注页数：{notes_count}{'（包含基础提示词）' if base_prompt else ''}")


def main():
    parser = argparse.ArgumentParser(description="合并幻灯片图片为 PPTX")
    parser.add_argument("dir", help="幻灯片目录")
    parser.add_argument("--output", "-o", help="输出文件名")
    args = parser.parse_args()

    slides = find_slide_images(args.dir)

    dir_name = Path(args.dir).name
    if dir_name == "slide-deck":
        dir_name = Path(args.dir).parent.name

    output = args.output or str(Path(args.dir) / f"{dir_name}.pptx")

    print(f"在 {args.dir} 中找到 {len(slides)} 张幻灯片\n")
    create_pptx(slides, output, args.dir)


if __name__ == "__main__":
    main()
